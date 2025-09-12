# Working AI-Powered Speaker Avatar for Presentations
# Fixed for Python 3.12 and Windows compatibility

import os
import sys
import subprocess
import logging
from pathlib import Path
from typing import List, Dict, Optional
import json
from datetime import datetime

# Core dependencies
import numpy as np
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

# Video processing
try:
    from moviepy.editor import (
        ImageClip, VideoFileClip, CompositeVideoClip, 
        concatenate_videoclips, AudioFileClip
    )
    MOVIEPY_AVAILABLE = True
except ImportError:
    MOVIEPY_AVAILABLE = False
    print("WARNING: MoviePy not available, video composition will be limited")

# Progress tracking
from tqdm import tqdm

# TTS Fallback
from tts_fallback import TTSManager

# Setup logging with Windows-compatible encoding
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('speaker_avatar.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class WorkingConfig:
    """Working configuration class"""
    
    def __init__(self):
        self.config = {
            "directories": {
                "ppt_input": "ppt_input",
                "face_images": "face_images",
                "audio": "audio",
                "video": "video",
                "output": "output",
                "slide_images": "slide_images",
                "temp": "temp"
            },
            "files": {
                "face_image": "face.png",
                "final_video": "final_presentation.mp4"
            },
            "wav2lip": {
                "path": "Wav2Lip",
                "checkpoint": "checkpoints/wav2lip_gan.pth"
            },
            "tts": {
                "model": "tts_models/en/ljspeech/tacotron2-DDC",
                "voice": "default",
                "speed": 1.0,
                "pitch": 1.0
            },
            "video": {
                "fps": 25,
                "resolution": [1920, 1080],
                "avatar_size_ratio": 0.25,
                "avatar_position": "bottom_right"
            },
            "processing": {
                "extract_slide_content": True,
                "extract_speaker_notes": True,
                "combine_text_sources": True,
                "skip_existing": True
            }
        }

class PowerPointProcessor:
    """PowerPoint text extraction"""
    
    def __init__(self, config: WorkingConfig):
        self.config = config
        self.logger = logging.getLogger(__name__)
    
    def extract_text_from_slide(self, slide) -> Dict[str, str]:
        """Extract text content from a single slide"""
        slide_data = {
            'title': '',
            'content': '',
            'speaker_notes': '',
            'combined_text': ''
        }
        
        # Extract title
        if slide.shapes.title:
            slide_data['title'] = slide.shapes.title.text.strip()
        
        # Extract content from text boxes
        content_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape != slide.shapes.title:  # Skip title
                    content_parts.append(shape.text.strip())
        
        slide_data['content'] = '\n'.join(content_parts)
        
        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_data['speaker_notes'] = slide.notes_slide.notes_text_frame.text.strip()
        
        # Combine text intelligently
        combined_parts = []
        if slide_data['title']:
            combined_parts.append(f"Title: {slide_data['title']}")
        if slide_data['content']:
            combined_parts.append(f"Content: {slide_data['content']}")
        if slide_data['speaker_notes']:
            combined_parts.append(f"Notes: {slide_data['speaker_notes']}")
        
        slide_data['combined_text'] = '\n\n'.join(combined_parts) if combined_parts else slide_data['speaker_notes'] or slide_data['content']
        
        return slide_data
    
    def generate_tutor_explanation(self, slide_data: Dict[str, str]) -> str:
        """Generate tutor-style explanation from slide content"""
        title = slide_data.get('title', '')
        content = slide_data.get('content', '')
        notes = slide_data.get('speaker_notes', '')
        
        # Create a natural, explanatory narration
        explanation_parts = []
        
        # Start with a greeting or transition
        if slide_data['slide_number'] == 1:
            explanation_parts.append("Welcome! Let me explain this presentation to you.")
        
        # Add title explanation
        if title:
            explanation_parts.append(f"On this slide, we're discussing {title.lower()}")
        
        # Process content to create explanations
        if content:
            # Split content into sentences and create explanations
            sentences = [s.strip() for s in content.split('.') if s.strip()]
            
            for sentence in sentences:
                if len(sentence) > 10:  # Only process substantial sentences
                    # Add explanatory phrases
                    if 'system' in sentence.lower():
                        explanation_parts.append(f"Let me explain: {sentence}")
                    elif 'leverages' in sentence.lower() or 'uses' in sentence.lower():
                        explanation_parts.append(f"Here's how it works: {sentence}")
                    elif 'can' in sentence.lower() or 'will' in sentence.lower():
                        explanation_parts.append(f"The key benefits are: {sentence}")
                    else:
                        explanation_parts.append(f"To elaborate: {sentence}")
        
        # Add notes if available
        if notes and len(notes) > 20:
            explanation_parts.append(f"Additional context: {notes}")
        
        # If no good content, fall back to original text
        if not explanation_parts:
            return slide_data['combined_text']
        
        # Join with natural transitions
        explanation = '. '.join(explanation_parts)
        if not explanation.endswith('.'):
            explanation += '.'
        
        return explanation
    
    def extract_slides_from_pptx(self, pptx_path: str) -> List[Dict[str, str]]:
        """Extract text from all slides in a PowerPoint presentation"""
        self.logger.info(f"Extracting text from {pptx_path}...")
        
        try:
            prs = Presentation(pptx_path)
            slides_data = []
            
            for i, slide in enumerate(prs.slides):
                slide_data = self.extract_text_from_slide(slide)
                slide_data['slide_number'] = i + 1
                
                # Generate tutor-style explanation after slide_number is set
                slide_data['narration_text'] = self.generate_tutor_explanation(slide_data)
                
                slides_data.append(slide_data)
                
                title_preview = slide_data.get('title', 'No title')[:50]
                self.logger.info(f"   -> Slide {i+1}: '{title_preview}...'")
            
            self.logger.info(f"SUCCESS: Extracted text from {len(slides_data)} slides")
            return slides_data
            
        except Exception as e:
            self.logger.error(f"ERROR: Reading PowerPoint file: {e}")
            return []
    
    def create_slide_image(self, slide, slide_number: int, output_path: str):
        """Create a simple slide image"""
        try:
            # Create image
            img = Image.new('RGB', (1920, 1080), 'white')
            draw = ImageDraw.Draw(img)
            
            # Try to load fonts
            try:
                title_font = ImageFont.truetype("arial.ttf", 48)
                content_font = ImageFont.truetype("arial.ttf", 32)
            except:
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
            
            # Draw slide number
            draw.text((1800, 20), f"Slide {slide_number}", fill='gray', font=content_font)
            
            # Draw title
            if slide.shapes.title and slide.shapes.title.text.strip():
                title_text = slide.shapes.title.text.strip()
                draw.text((50, 50), title_text, fill='black', font=title_font)
            
            # Draw content
            y_pos = 150
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip() and shape != slide.shapes.title:
                    text = shape.text.strip()
                    # Simple text wrapping
                    words = text.split()
                    line = ""
                    for word in words:
                        test_line = line + word + " "
                        bbox = draw.textbbox((0, 0), test_line, font=content_font)
                        if bbox[2] - bbox[0] > 1800:  # If line too long
                            if line:
                                draw.text((50, y_pos), line.strip(), fill='black', font=content_font)
                                y_pos += 40
                                line = word + " "
                        else:
                            line = test_line
                    if line:
                        draw.text((50, y_pos), line.strip(), fill='black', font=content_font)
                        y_pos += 40
            
            img.save(output_path)
            self.logger.debug(f"Created slide image: {output_path}")
            
        except Exception as e:
            self.logger.error(f"Error creating slide image: {e}")

class TTSProcessor:
    """TTS processor with fallback support"""
    
    def __init__(self, config: WorkingConfig):
        self.config = config
        self.logger = logging.getLogger(__name__)
        self.tts_manager = TTSManager(self.config.config['tts'])
    
    def generate_audio(self, text: str, output_path: str) -> bool:
        """Generate audio from text"""
        if not self.tts_manager.is_available():
            self.logger.error("No TTS system available")
            return False
        
        try:
            success = self.tts_manager.generate_audio(text, output_path)
            if success:
                self.logger.info(f"SUCCESS: Generated audio: {output_path}")
            return success
        except Exception as e:
            self.logger.error(f"ERROR: Generating audio: {e}")
            return False
    
    def generate_audio_batch(self, slides_data: List[Dict], output_dir: str) -> List[str]:
        """Generate audio for all slides"""
        self.logger.info("Generating audio for all slides...")
        
        audio_files = []
        for slide_data in tqdm(slides_data, desc="Generating audio"):
            slide_num = slide_data['slide_number']
            text = slide_data['narration_text']  # Use tutor-style explanation
            
            if not text.strip():
                self.logger.warning(f"WARNING: No text found for slide {slide_num}, skipping")
                continue
            
            output_path = os.path.join(output_dir, f"slide_{slide_num:03d}.wav")
            
            # Skip if file exists and skip_existing is enabled
            if (os.path.exists(output_path) and 
                self.config.config['processing']['skip_existing']):
                self.logger.info(f"SKIP: Slide {slide_num}, audio already exists")
                audio_files.append(output_path)
                continue
            
            if self.generate_audio(text, output_path):
                audio_files.append(output_path)
        
        self.logger.info(f"SUCCESS: Generated {len(audio_files)} audio files")
        return audio_files

class Wav2LipProcessor:
    """Wav2Lip processor"""
    
    def __init__(self, config: WorkingConfig):
        self.config = config
        self.logger = logging.getLogger(__name__)
        self.wav2lip_path = os.path.join(config.config['wav2lip']['path'])
        self.checkpoint_path = os.path.join(self.wav2lip_path, config.config['wav2lip']['checkpoint'])
    
    def animate_face(self, audio_path: str, face_path: str, output_path: str) -> bool:
        """Animate face using Wav2Lip"""
        try:
            command = [
                "python", os.path.join(self.wav2lip_path, "inference.py"),
                "--checkpoint_path", self.checkpoint_path,
                "--face", face_path,
                "--audio", audio_path,
                "--outfile", output_path,
                "--static", "True",
                "--fps", str(self.config.config['video']['fps']),
                "--pads", "0", "20", "0", "0",  # Add padding to include chin and mouth area
                "--resize_factor", "1",  # Don't resize down
                "--wav2lip_batch_size", "4"  # Smaller batch size for stability
            ]
            
            self.logger.info(f"Animating face for {os.path.basename(audio_path)}...")
            
            result = subprocess.run(
                command, 
                check=True, 
                capture_output=True, 
                text=True,
                cwd=os.getcwd()
            )
            
            self.logger.info(f"SUCCESS: Animation complete: {output_path}")
            return True
            
        except subprocess.CalledProcessError as e:
            self.logger.error(f"ERROR: Wav2Lip error: {e.stderr}")
            return False
        except Exception as e:
            self.logger.error(f"ERROR: Face animation: {e}")
            return False
    
    def animate_faces_batch(self, audio_files: List[str], face_path: str, output_dir: str) -> List[str]:
        """Animate faces for all audio files"""
        self.logger.info("Starting batch face animation...")
        
        animated_videos = []
        for audio_path in tqdm(audio_files, desc="Animating faces"):
            base_name = os.path.splitext(os.path.basename(audio_path))[0]
            output_path = os.path.join(output_dir, f"{base_name}_animated.mp4")
            
            if (os.path.exists(output_path) and 
                self.config.config['processing']['skip_existing']):
                self.logger.info(f"SKIP: {base_name}, video already exists")
                animated_videos.append(output_path)
                continue
            
            if self.animate_face(audio_path, face_path, output_path):
                animated_videos.append(output_path)
        
        self.logger.info(f"SUCCESS: Generated {len(animated_videos)} animated videos")
        return animated_videos

class VideoComposer:
    """Video composer"""
    
    def __init__(self, config: WorkingConfig):
        self.config = config
        self.logger = logging.getLogger(__name__)
    
    def compose_final_video(self, animated_videos: List[str], slide_images: List[str], 
                          output_path: str) -> bool:
        """Compose final video with slides and avatar"""
        self.logger.info("Composing final video...")
        
        if not MOVIEPY_AVAILABLE:
            self.logger.error("ERROR: MoviePy not available, cannot compose video")
            return False
        
        try:
            final_clips = []
            video_config = self.config.config['video']
            
            for i, (animated_video, slide_image) in enumerate(zip(animated_videos, slide_images)):
                if not os.path.exists(animated_video) or not os.path.exists(slide_image):
                    self.logger.warning(f"WARNING: Missing files for slide {i+1}, skipping")
                    continue
                
                # Load video and image clips
                video_clip = VideoFileClip(animated_video)
                slide_clip = ImageClip(slide_image).set_duration(video_clip.duration)
                
                # Resize avatar
                avatar_ratio = video_config['avatar_size_ratio']
                avatar_height = int(slide_clip.h * avatar_ratio)
                video_clip_resized = video_clip.resize(height=avatar_height)
                
                # Position avatar
                position = self._get_avatar_position(
                    video_config['avatar_position'], 
                    slide_clip.size, 
                    video_clip_resized.size
                )
                video_clip_resized = video_clip_resized.set_position(position)
                
                # Composite the clips
                final_clip = CompositeVideoClip([slide_clip, video_clip_resized])
                final_clips.append(final_clip)
            
            if not final_clips:
                self.logger.error("ERROR: No clips to compose")
                return False
            
            # Concatenate all clips
            final_video = concatenate_videoclips(final_clips, method="compose")
            
            # Write final video
            final_video.write_videofile(
                output_path,
                codec="libx264",
                audio_codec="aac",
                fps=video_config['fps']
            )
            
            self.logger.info(f"SUCCESS: Final video saved: {output_path}")
            return True
            
        except Exception as e:
            self.logger.error(f"ERROR: Composing final video: {e}")
            return False
    
    def _get_avatar_position(self, position: str, slide_size: tuple, avatar_size: tuple) -> str:
        """Get avatar position based on configuration"""
        if position == "bottom_right":
            return ("right", "bottom")
        elif position == "bottom_left":
            return ("left", "bottom")
        elif position == "top_right":
            return ("right", "top")
        elif position == "top_left":
            return ("left", "top")
        else:
            return ("right", "bottom")  # Default

class WorkingSpeakerAvatarSystem:
    """Working main system class"""
    
    def __init__(self):
        self.config = WorkingConfig()
        self.logger = logging.getLogger(__name__)
        
        # Initialize processors
        self.ppt_processor = PowerPointProcessor(self.config)
        self.tts_processor = TTSProcessor(self.config)
        self.wav2lip_processor = Wav2LipProcessor(self.config)
        self.video_composer = VideoComposer(self.config)
        
        # Create directories
        self._create_directories()
    
    def _create_directories(self):
        """Create all necessary directories"""
        for dir_name in self.config.config['directories'].values():
            os.makedirs(dir_name, exist_ok=True)
        self.logger.info("SUCCESS: All directories created")
    
    def process_presentation(self, pptx_path: str) -> bool:
        """Process a single PowerPoint presentation"""
        self.logger.info(f"Processing presentation: {pptx_path}")
        
        # Step 1: Extract text from PowerPoint
        slides_data = self.ppt_processor.extract_slides_from_pptx(pptx_path)
        if not slides_data:
            return False
        
        # Step 2: Create slide images
        slide_images = []
        for i, slide_data in enumerate(slides_data):
            slide_image_path = os.path.join(
                self.config.config['directories']['slide_images'],
                f"slide_{i+1:03d}.png"
            )
            # For now, create a simple placeholder
            # In a full implementation, you'd render the actual slide
            self.ppt_processor.create_slide_image(
                Presentation(pptx_path).slides[i], 
                i + 1, 
                slide_image_path
            )
            slide_images.append(slide_image_path)
        
        # Step 3: Generate audio
        audio_files = self.tts_processor.generate_audio_batch(
            slides_data, 
            self.config.config['directories']['audio']
        )
        if not audio_files:
            return False
        
        # Step 4: Animate faces
        face_path = os.path.join(
            self.config.config['directories']['face_images'],
            self.config.config['files']['face_image']
        )
        
        animated_videos = self.wav2lip_processor.animate_faces_batch(
            audio_files,
            face_path,
            self.config.config['directories']['video']
        )
        if not animated_videos:
            return False
        
        # Step 5: Compose final video
        output_path = os.path.join(
            self.config.config['directories']['output'],
            self.config.config['files']['final_video']
        )
        
        success = self.video_composer.compose_final_video(
            animated_videos, 
            slide_images, 
            output_path
        )
        
        if success:
            self.logger.info("SUCCESS: Presentation processing completed!")
        
        return success

def main():
    """Main function"""
    print("AI-Powered Speaker Avatar for Presentations (Working Version)")
    print("=" * 60)
    
    # Initialize system
    system = WorkingSpeakerAvatarSystem()
    
    # Check if we have a presentation to process
    ppt_input_dir = system.config.config['directories']['ppt_input']
    pptx_files = [f for f in os.listdir(ppt_input_dir) if f.endswith('.pptx')]
    
    if not pptx_files:
        print(f"ERROR: No PowerPoint files found in {ppt_input_dir}")
        print("Please place your .pptx files in the ppt_input directory")
        return
    
    # Process presentations
    if len(pptx_files) == 1:
        # Single presentation
        pptx_path = os.path.join(ppt_input_dir, pptx_files[0])
        success = system.process_presentation(pptx_path)
        if success:
            print("SUCCESS: Single presentation processed successfully!")
        else:
            print("ERROR: Failed to process presentation")
    else:
        # Multiple presentations
        for pptx_file in pptx_files:
            pptx_path = os.path.join(ppt_input_dir, pptx_file)
            print(f"Processing {pptx_file}...")
            success = system.process_presentation(pptx_path)
            if success:
                print(f"SUCCESS: {pptx_file} processed successfully!")
            else:
                print(f"ERROR: Failed to process {pptx_file}")

if __name__ == "__main__":
    main()
