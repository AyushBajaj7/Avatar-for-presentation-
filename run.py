#!/usr/bin/env python3
"""
Easy AI Speaker Avatar System - Single Command Runner
Run with: python run.py
"""

import os
import sys
import subprocess
import logging
import time
import json
import shutil
from pathlib import Path
from typing import List, Dict, Optional
from datetime import datetime
import tempfile
import threading
import base64
from io import BytesIO

# Check if dependencies are available
def check_dependencies():
    """Check if all required dependencies are available"""
    required_packages = {
        "pptx": "python-pptx",
        "PIL": "Pillow", 
        "cv2": "opencv-python",
        "pyttsx3": "pyttsx3",
        "tqdm": "tqdm",
        "numpy": "numpy",
        "flask": "flask",
        "flask_cors": "flask-cors",
        "google.generativeai": "google-generativeai",
        "PyPDF2": "PyPDF2",
        "pdf2image": "pdf2image"
    }
    
    missing_packages = []
    
    for import_name, package_name in required_packages.items():
        try:
            __import__(import_name)
        except ImportError:
            missing_packages.append(package_name)
    
    if missing_packages:
        print("❌ Missing required packages:")
        for package in missing_packages:
            print(f"  - {package}")
        print()
        print("🔧 Please run 'python download.py' first to install all dependencies.")
        print("   Or install manually: pip install " + " ".join(missing_packages))
        sys.exit(1)
    
    print("✅ All dependencies are available")

# Check dependencies before importing
check_dependencies()

# Now import the dependencies
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import numpy as np
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
from tqdm import tqdm

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('easy_system.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Flask app setup
app = Flask(__name__)
CORS(app)

class EasyConfig:
    """Easy configuration class"""
    
    def __init__(self):
        self.config = {
            "directories": {
                "uploads": "uploads",
                "audio": "audio",
                "video": "video",
                "output": "output",
                "slide_images": "slide_images",
                "temp": "temp"
            },
            "files": {
                "final_video": "final_presentation_with_slides.mp4"
            },
            "wav2lip": {
                "path": "Wav2Lip",
                "checkpoint": "checkpoints/wav2lip_gan.pth"
            },
            "tts": {
                "model": "tts_models/en/ljspeech/tacotron2-DDC",
                "voice": "default",
                "speed": 1.0,
                "pitch": 1.0
            },
            "video": {
                "fps": 25,
                "resolution": [1920, 1080],
                "avatar_size_ratio": 0.25,
                "avatar_position": "bottom_right"
            },
            "processing": {
                "extract_slide_content": True,
                "extract_speaker_notes": True,
                "combine_text_sources": True,
                "skip_existing": True
            },
            "gemini": {
                "api_key": "",
                "enabled": False
            }
        }

class TTSManager:
    """Simple TTS Manager with pyttsx3 fallback"""
    
    def __init__(self, config):
        self.config = config
        self.engine = None
        self._init_engine()
    
    def _init_engine(self):
        """Initialize TTS engine"""
        try:
            import pyttsx3
            self.engine = pyttsx3.init()
            self.engine.setProperty('rate', 150)
            self.engine.setProperty('volume', 0.9)
            logger.info("SUCCESS: TTS engine initialized with pyttsx3")
        except Exception as e:
            logger.error(f"ERROR: Failed to initialize TTS: {e}")
            self.engine = None
    
    def reset_engine(self):
        """Reset the TTS engine completely"""
        try:
            if self.engine:
                self.engine.stop()
            time.sleep(1)
            self._init_engine()
            logger.info("SUCCESS: TTS engine reset")
        except Exception as e:
            logger.error(f"ERROR: Failed to reset TTS engine: {e}")
            self.engine = None
    
    def set_voice(self, voice_name: str):
        """Set the TTS voice"""
        if not self.is_available():
            return False
        
        try:
            voices = self.engine.getProperty('voices')
            for voice in voices:
                if voice.name == voice_name:
                    self.engine.setProperty('voice', voice.id)
                    logger.info(f"SUCCESS: Voice set to {voice_name}")
                    return True
            
            logger.warning(f"WARNING: Voice '{voice_name}' not found, using default")
            return False
        except Exception as e:
            logger.error(f"ERROR: Failed to set voice: {e}")
            return False
    
    def is_available(self):
        """Check if TTS is available"""
        return self.engine is not None
    
    def generate_audio(self, text: str, output_path: str) -> bool:
        """Generate audio from text with timeout"""
        if not self.is_available():
            logger.error("TTS engine not available")
            return False
        
        try:
            # Stop any running operations
            try:
                self.engine.stop()
            except:
                pass
            
            # Wait a moment for cleanup
            time.sleep(0.5)
            
            # Generate audio with timeout
            import threading
            import queue
            
            result_queue = queue.Queue()
            error_queue = queue.Queue()
            
            def generate_audio_thread():
                try:
                    self.engine.save_to_file(text, output_path)
                    self.engine.runAndWait()
                    result_queue.put(True)
                except Exception as e:
                    error_queue.put(e)
            
            # Start TTS generation in a separate thread
            tts_thread = threading.Thread(target=generate_audio_thread)
            tts_thread.daemon = True
            tts_thread.start()
            
            # Wait for completion with timeout
            tts_thread.join(timeout=30)  # 30 second timeout
            
            if tts_thread.is_alive():
                logger.error("ERROR: TTS generation timed out after 30 seconds")
                # Force stop the engine
                try:
                    self.engine.stop()
                except:
                    pass
                return False
            
            # Check for errors
            if not error_queue.empty():
                error = error_queue.get()
                logger.error(f"ERROR: TTS generation failed: {error}")
                return False
            
            # Wait for file to be written
            time.sleep(1)
            
            # Verify file was created and has content
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                logger.info(f"SUCCESS: Generated audio file: {output_path}")
                return True
            else:
                logger.error(f"ERROR: Audio file not created or empty: {output_path}")
                return False
                
        except Exception as e:
            logger.error(f"ERROR: TTS generation failed: {e}")
            # Reset the engine completely
            try:
                self.engine.stop()
                time.sleep(1)
                self._init_engine()
            except Exception as reset_error:
                logger.error(f"ERROR: Failed to reset TTS engine: {reset_error}")
            return False

class PowerPointProcessor:
    """PowerPoint text extraction with Gemini enhancement"""
    
    def __init__(self, config: EasyConfig):
        self.config = config
        self.logger = logging.getLogger(__name__)
    
    def extract_text_from_slide(self, slide) -> Dict[str, str]:
        """Extract text content from a single slide"""
        slide_data = {
            'title': '',
            'content': '',
            'speaker_notes': '',
            'combined_text': ''
        }
        
        # Extract title
        if slide.shapes.title:
            slide_data['title'] = slide.shapes.title.text.strip()
        
        # Extract content from text boxes
        content_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape != slide.shapes.title:  # Skip title
                    content_parts.append(shape.text.strip())
        
        slide_data['content'] = '\n'.join(content_parts)
        
        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_data['speaker_notes'] = slide.notes_slide.notes_text_frame.text.strip()
        
        # Combine text intelligently
        combined_parts = []
        if slide_data['title']:
            combined_parts.append(f"Title: {slide_data['title']}")
        if slide_data['content']:
            combined_parts.append(f"Content: {slide_data['content']}")
        if slide_data['speaker_notes']:
            combined_parts.append(f"Notes: {slide_data['speaker_notes']}")
        
        slide_data['combined_text'] = '\n\n'.join(combined_parts) if combined_parts else slide_data['speaker_notes'] or slide_data['content']
        
        return slide_data
    
    def enhance_with_gemini(self, slide_data: Dict[str, str]) -> str:
        """Enhance slide content using Gemini API"""
        if not self.config.config['gemini']['enabled'] or not self.config.config['gemini']['api_key']:
            return self.generate_tutor_explanation(slide_data)
        
        try:
            import google.generativeai as genai
            
            genai.configure(api_key=self.config.config['gemini']['api_key'])
            model = genai.GenerativeModel('gemini-pro')
            
            prompt = f"""
            You are an expert presentation tutor. Create a natural, engaging explanation for this slide content:
            
            Title: {slide_data.get('title', '')}
            Content: {slide_data.get('content', '')}
            Speaker Notes: {slide_data.get('speaker_notes', '')}
            
            Please create a conversational, tutor-style explanation that:
            1. Explains the concepts clearly
            2. Uses natural transitions
            3. Is engaging and educational
            4. Flows naturally for speech
            
            Keep it concise but informative. Start with a brief introduction if this is slide 1.
            """
            
            response = model.generate_content(prompt)
            return response.text.strip()
            
        except Exception as e:
            self.logger.warning(f"Gemini enhancement failed: {e}, using fallback")
            return self.generate_tutor_explanation(slide_data)
    
    def generate_tutor_explanation(self, slide_data: Dict[str, str]) -> str:
        """Generate tutor-style explanation from slide content"""
        title = slide_data.get('title', '')
        content = slide_data.get('content', '')
        notes = slide_data.get('speaker_notes', '')
        
        # Create a natural, explanatory narration
        explanation_parts = []
        
        # Start with a greeting or transition
        if slide_data.get('slide_number', 1) == 1:
            explanation_parts.append("Welcome! Let me explain this presentation to you.")
        
        # Add title explanation
        if title:
            explanation_parts.append(f"On this slide, we're discussing {title.lower()}")
        
        # Process content to create explanations
        if content:
            # Split content into sentences and create explanations
            sentences = [s.strip() for s in content.split('.') if s.strip()]
            
            for sentence in sentences:
                if len(sentence) > 10:  # Only process substantial sentences
                    # Add explanatory phrases
                    if 'system' in sentence.lower():
                        explanation_parts.append(f"Let me explain: {sentence}")
                    elif 'leverages' in sentence.lower() or 'uses' in sentence.lower():
                        explanation_parts.append(f"Here's how it works: {sentence}")
                    elif 'can' in sentence.lower() or 'will' in sentence.lower():
                        explanation_parts.append(f"The key benefits are: {sentence}")
                    else:
                        explanation_parts.append(f"To elaborate: {sentence}")
        
        # Add notes if available
        if notes and len(notes) > 20:
            explanation_parts.append(f"Additional context: {notes}")
        
        # If no good content, fall back to original text
        if not explanation_parts:
            return slide_data['combined_text']
        
        # Join with natural transitions
        explanation = '. '.join(explanation_parts)
        if not explanation.endswith('.'):
            explanation += '.'
        
        return explanation
    
    def extract_slides_from_pptx(self, pptx_path: str) -> List[Dict[str, str]]:
        """Extract text from all slides in a PowerPoint presentation"""
        self.logger.info(f"Extracting text from {pptx_path}...")
        
        try:
            prs = Presentation(pptx_path)
            slides_data = []
            
            for i, slide in enumerate(prs.slides):
                slide_data = self.extract_text_from_slide(slide)
                slide_data['slide_number'] = i + 1
                
                # Generate enhanced explanation
                slide_data['narration_text'] = self.enhance_with_gemini(slide_data)
                
                slides_data.append(slide_data)
                
                title_preview = slide_data.get('title', 'No title')[:50]
                self.logger.info(f"   -> Slide {i+1}: '{title_preview}...'")
            
            self.logger.info(f"SUCCESS: Extracted text from {len(slides_data)} slides")
            return slides_data
            
        except Exception as e:
            self.logger.error(f"ERROR: Reading PowerPoint file: {e}")
            return []
    
    def extract_slides_from_pdf(self, pdf_path: str) -> List[Dict[str, str]]:
        """Extract text from all pages in a PDF presentation"""
        self.logger.info(f"Extracting text from {pdf_path}...")
        
        try:
            import PyPDF2
            
            slides_data = []
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for i, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    
                    slide_data = {
                        'title': f'Page {i + 1}',
                        'content': text,
                        'speaker_notes': '',
                        'combined_text': text,
                        'slide_number': i + 1
                    }
                    
                    # Generate enhanced explanation
                    slide_data['narration_text'] = self.enhance_with_gemini(slide_data)
                    
                    slides_data.append(slide_data)
                    
                    content_preview = text[:50] if text else 'No content'
                    self.logger.info(f"   -> Page {i+1}: '{content_preview}...'")
            
            self.logger.info(f"SUCCESS: Extracted text from {len(slides_data)} pages")
            return slides_data
            
        except Exception as e:
            self.logger.error(f"ERROR: Reading PDF file: {e}")
            return []
    
    def extract_slides_from_file(self, file_path: str) -> List[Dict[str, str]]:
        """Extract slides from either PPTX or PDF file"""
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pptx':
            return self.extract_slides_from_pptx(file_path)
        elif file_ext == '.pdf':
            return self.extract_slides_from_pdf(file_path)
        else:
            self.logger.error(f"ERROR: Unsupported file format: {file_ext}")
            return []

class SlideRenderer:
    """Enhanced slide renderer for PowerPoint presentations"""
    
    def __init__(self, width: int = 1920, height: int = 1080):
        self.width = width
        self.height = height
        self.logger = logging.getLogger(__name__)
        
        # Try to load a system font
        try:
            self.title_font = ImageFont.truetype("arial.ttf", 48)
            self.content_font = ImageFont.truetype("arial.ttf", 32)
            self.notes_font = ImageFont.truetype("arial.ttf", 24)
        except:
            # Fallback to default font
            self.title_font = ImageFont.load_default()
            self.content_font = ImageFont.load_default()
            self.notes_font = ImageFont.load_default()
    
    def render_slide(self, slide, slide_number: int) -> Image.Image:
        """Render a single slide as an image"""
        # Create a new image with white background
        img = Image.new('RGB', (self.width, self.height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Define layout areas
        margin = 50
        title_area = (margin, margin, self.width - margin, 150)
        content_area = (margin, 180, self.width - margin, self.height - 200)
        notes_area = (margin, self.height - 150, self.width - margin, self.height - margin)
        
        # Draw slide number
        draw.text((self.width - 100, 20), f"Slide {slide_number}", 
                 fill='gray', font=self.content_font)
        
        # Extract and render title
        if slide.shapes.title and slide.shapes.title.text.strip():
            title_text = slide.shapes.title.text.strip()
            self._draw_wrapped_text(draw, title_text, title_area, self.title_font, 'black')
        
        # Extract and render content
        content_text = self._extract_slide_content(slide)
        if content_text:
            self._draw_wrapped_text(draw, content_text, content_area, self.content_font, 'black')
        
        # Extract and render speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                # Draw a separator line
                draw.line([margin, self.height - 160, self.width - margin, self.height - 160], 
                         fill='lightgray', width=2)
                self._draw_wrapped_text(draw, f"Notes: {notes_text}", notes_area, 
                                      self.notes_font, 'darkgray')
        
        return img
    
    def _extract_slide_content(self, slide) -> str:
        """Extract all text content from slide shapes"""
        content_parts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Skip the title shape
                if shape != slide.shapes.title:
                    content_parts.append(shape.text.strip())
        
        return '\n'.join(content_parts)
    
    def _draw_wrapped_text(self, draw, text: str, area: tuple, font, color: str):
        """Draw text with word wrapping within the specified area"""
        x, y, max_x, max_y = area
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            test_line = ' '.join(current_line + [word])
            bbox = draw.textbbox((0, 0), test_line, font=font)
            text_width = bbox[2] - bbox[0]
            
            if text_width <= max_x - x:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                    current_line = [word]
                else:
                    # Single word is too long, add it anyway
                    lines.append(word)
        
        if current_line:
            lines.append(' '.join(current_line))
        
        # Draw the lines
        line_height = 40  # Approximate line height
        for i, line in enumerate(lines):
            line_y = y + (i * line_height)
            if line_y + line_height <= max_y:
                draw.text((x, line_y), line, fill=color, font=font)
    
    def render_presentation(self, file_path: str, output_dir: str) -> List[str]:
        """Render all slides from a presentation (PPTX or PDF)"""
        self.logger.info(f"Rendering slides from {file_path}")
        
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            slide_images = []
            
            if file_ext == '.pptx':
                # Handle PowerPoint files
                prs = Presentation(file_path)
                
                for i, slide in enumerate(prs.slides):
                    slide_image = self.render_slide(slide, i + 1)
                    output_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
                    slide_image.save(output_path)
                    slide_images.append(output_path)
                    self.logger.info(f"   -> Rendered slide {i+1}")
                    
            elif file_ext == '.pdf':
                # Handle PDF files
                from pdf2image import convert_from_path
                
                pages = convert_from_path(file_path, dpi=300)
                
                for i, page in enumerate(pages):
                    # Resize to our standard resolution
                    page = page.resize((self.width, self.height), Image.Resampling.LANCZOS)
                    output_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
                    page.save(output_path)
                    slide_images.append(output_path)
                    self.logger.info(f"   -> Rendered page {i+1}")
            
            self.logger.info(f"SUCCESS: Rendered {len(slide_images)} slides")
            return slide_images
            
        except Exception as e:
            self.logger.error(f"ERROR: Rendering slides: {e}")
            return []

class TTSProcessor:
    """TTS processor with fallback support"""
    
    def __init__(self, config: EasyConfig):
        self.config = config
        self.logger = logging.getLogger(__name__)
        self.tts_manager = TTSManager(self.config.config['tts'])
        self.system_instance = None  # Will be set by the main system
    
    def generate_audio(self, text: str, output_path: str) -> bool:
        """Generate audio from text"""
        if not self.tts_manager.is_available():
            self.logger.error("No TTS system available")
            return False
        
        try:
            success = self.tts_manager.generate_audio(text, output_path)
            if success:
                self.logger.info(f"SUCCESS: Generated audio: {output_path}")
            return success
        except Exception as e:
            self.logger.error(f"ERROR: Generating audio: {e}")
            return False
    
    def set_system_instance(self, system_instance):
        """Set reference to main system for status updates"""
        self.system_instance = system_instance
    
    def update_status(self, progress: int, current_step: str):
        """Update the global processing status"""
        if self.system_instance:
            self.system_instance.processing_status['progress'] = progress
            self.system_instance.processing_status['current_step'] = current_step
            # Show status update prominently in terminal
            print(f"\n🔄 STATUS UPDATE: {progress}% - {current_step}")
            self.logger.info(f"Status updated: {progress}% - {current_step}")
    
    def generate_audio_batch(self, slides_data: List[Dict], output_dir: str) -> List[str]:
        """Generate audio for all slides"""
        self.logger.info("Generating audio for all slides...")
        
        # Ensure output directory exists
        os.makedirs(output_dir, exist_ok=True)
        
        audio_files = []
        failed_slides = []
        total_slides = len(slides_data)
        
        for i, slide_data in enumerate(slides_data):
            slide_num = slide_data['slide_number']
            text = slide_data['narration_text']  # Use enhanced explanation
            
            # Update progress and status
            progress = 25 + (i * 15) // total_slides  # 25-40% range
            current_step = f"Generating audio for slide {slide_num}/{total_slides}"
            print(f"\n🎵 PROCESSING SLIDE {slide_num}/{total_slides} - Generating audio...")
            self.logger.info(f"Processing slide {slide_num}/{total_slides} (Progress: {progress}%)")
            
            # Update global status for frontend
            self.update_status(progress, current_step)
            
            if not text.strip():
                self.logger.warning(f"WARNING: No text found for slide {slide_num}, skipping")
                failed_slides.append(slide_num)
                continue
            
            output_path = os.path.join(output_dir, f"slide_{slide_num:03d}.wav")
            
            # Skip if file exists and skip_existing is enabled
            if (os.path.exists(output_path) and 
                self.config.config['processing']['skip_existing'] and
                os.path.getsize(output_path) > 0):
                self.logger.info(f"SKIP: Slide {slide_num}, audio already exists")
                audio_files.append(output_path)
                continue
            
            # Try to generate audio with retry logic
            max_retries = 3
            success = False
            
            for attempt in range(max_retries):
                current_step = f"Generating audio for slide {slide_num}/{total_slides} (attempt {attempt + 1}/{max_retries})"
                print(f"🔄 ATTEMPT {attempt + 1}/{max_retries} for slide {slide_num}")
                self.logger.info(f"Generating audio for slide {slide_num} (attempt {attempt + 1}/{max_retries})")
                self.update_status(progress, current_step)
                
                if self.generate_audio(text, output_path):
                    audio_files.append(output_path)
                    success = True
                    print(f"✅ AUDIO GENERATED FOR SLIDE {slide_num}/{total_slides} - Moving to next slide")
                    self.logger.info(f"SUCCESS: Audio generated for slide {slide_num}")
                    # Update status to show completion
                    current_step = f"Audio generated for slide {slide_num}/{total_slides} - Moving to next slide"
                    self.update_status(progress, current_step)
                    break
                else:
                    print(f"❌ ATTEMPT {attempt + 1} FAILED for slide {slide_num}")
                    self.logger.warning(f"Attempt {attempt + 1} failed for slide {slide_num}")
                    if attempt < max_retries - 1:
                        # Reset engine between retries
                        print(f"🔄 Resetting TTS engine before retry {attempt + 2}")
                        self.logger.info(f"Resetting TTS engine before retry {attempt + 2}")
                        self.tts_manager.reset_engine()
                        time.sleep(2)  # Wait before retry
            
            if not success:
                print(f"❌ FAILED TO GENERATE AUDIO FOR SLIDE {slide_num} after {max_retries} attempts - SKIPPING")
                self.logger.error(f"ERROR: Failed to generate audio for slide {slide_num} after {max_retries} attempts")
                failed_slides.append(slide_num)
            else:
                # Reset engine after each successful generation to prevent sticking
                self.logger.info(f"Resetting TTS engine after slide {slide_num}")
                self.tts_manager.reset_engine()
                # Small delay to ensure status update is sent
                time.sleep(0.5)
        
        if failed_slides:
            self.logger.warning(f"WARNING: Failed to generate audio for slides: {failed_slides}")
        
        self.logger.info(f"SUCCESS: Generated {len(audio_files)} audio files out of {len(slides_data)} slides")
        return audio_files

class Wav2LipProcessor:
    """Wav2Lip processor"""
    
    def __init__(self, config: EasyConfig):
        self.config = config
        self.logger = logging.getLogger(__name__)
        self.wav2lip_path = os.path.join(config.config['wav2lip']['path'])
        self.checkpoint_path = os.path.join(self.wav2lip_path, config.config['wav2lip']['checkpoint'])
        self.system_instance = None  # Will be set by the main system
    
    def set_system_instance(self, system_instance):
        """Set reference to main system for status updates"""
        self.system_instance = system_instance
    
    def update_status(self, progress: int, current_step: str):
        """Update the global processing status"""
        if self.system_instance:
            self.system_instance.processing_status['progress'] = progress
            self.system_instance.processing_status['current_step'] = current_step
            # Show status update prominently in terminal
            print(f"\n🔄 STATUS UPDATE: {progress}% - {current_step}")
            self.logger.info(f"Status updated: {progress}% - {current_step}")
    
    def animate_face(self, audio_path: str, face_path: str, output_path: str) -> bool:
        """Animate face using Wav2Lip"""
        try:
            # Validate input files
            if not os.path.exists(audio_path):
                self.logger.error(f"ERROR: Audio file not found: {audio_path}")
                return False
            
            if not os.path.exists(face_path):
                self.logger.error(f"ERROR: Face image not found: {face_path}")
                return False
            
            # Check if Wav2Lip directory exists
            if not os.path.exists(self.wav2lip_path):
                self.logger.warning(f"WARNING: Wav2Lip directory not found: {self.wav2lip_path}")
                self.logger.warning("Creating a simple video with static face image instead...")
                return self._create_static_video(audio_path, face_path, output_path)
            
            # Check if checkpoint exists
            if not os.path.exists(self.checkpoint_path):
                self.logger.warning(f"WARNING: Wav2Lip checkpoint not found: {self.checkpoint_path}")
                self.logger.warning("Creating a simple video with static face image instead...")
                return self._create_static_video(audio_path, face_path, output_path)
            
            command = [
                "python", os.path.join(self.wav2lip_path, "inference.py"),
                "--checkpoint_path", self.checkpoint_path,
                "--face", face_path,
                "--audio", audio_path,
                "--outfile", output_path,
                "--static", "True",
                "--fps", str(self.config.config['video']['fps']),
                "--pads", "0", "20", "0", "0",  # Add padding to include chin and mouth area
                "--resize_factor", "1",  # Don't resize down
                "--wav2lip_batch_size", "2"  # Even smaller batch size for stability
            ]
            
            self.logger.info(f"Animating face for {os.path.basename(audio_path)}...")
            self.logger.info(f"Command: {' '.join(command)}")
            
            result = subprocess.run(
                command, 
                check=True, 
                capture_output=True, 
                text=True,
                cwd=os.getcwd(),
                timeout=300  # 5 minute timeout
            )
            
            # Verify output file was created
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                self.logger.info(f"SUCCESS: Animation complete: {output_path}")
                return True
            else:
                self.logger.error(f"ERROR: Animation output file not created or empty: {output_path}")
                return False
            
        except subprocess.TimeoutExpired:
            self.logger.error(f"ERROR: Wav2Lip timeout for {audio_path}")
            return False
        except subprocess.CalledProcessError as e:
            self.logger.error(f"ERROR: Wav2Lip error: {e.stderr}")
            self.logger.error(f"Command that failed: {' '.join(command)}")
            return False
        except Exception as e:
            self.logger.error(f"ERROR: Face animation: {e}")
            return False
    
    def _create_static_video(self, audio_path: str, face_path: str, output_path: str) -> bool:
        """Create a simple video with static face image and audio"""
        try:
            self.logger.info(f"Creating static video for {os.path.basename(audio_path)}...")
            
            # Get audio duration
            cmd_duration = [
                "ffprobe", "-v", "quiet", "-show_entries", "format=duration",
                "-of", "csv=p=0", audio_path
            ]
            
            result = subprocess.run(cmd_duration, capture_output=True, text=True)
            duration = float(result.stdout.strip()) if result.returncode == 0 else 5.0
            
            # Create video with static face image and audio
            # Use scale filter to ensure even dimensions for H.264 compatibility
            cmd = [
                "ffmpeg", "-y",
                "-loop", "1", "-i", face_path,  # Static face image
                "-i", audio_path,  # Audio file
                "-vf", "scale=426:640",  # Ensure even dimensions (427->426, 640 is already even)
                "-c:v", "libx264",
                "-c:a", "aac",
                "-t", str(duration),
                "-pix_fmt", "yuv420p",
                "-shortest",
                output_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0 and os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                self.logger.info(f"SUCCESS: Created static video: {output_path}")
                return True
            else:
                self.logger.error(f"ERROR: Failed to create static video: {result.stderr}")
                return False
                
        except Exception as e:
            self.logger.error(f"ERROR: Static video creation: {e}")
            return False
    
    def animate_faces_batch(self, audio_files: List[str], face_path: str, output_dir: str) -> List[str]:
        """Animate faces for all audio files"""
        self.logger.info("Starting batch face animation...")
        
        # Ensure output directory exists
        os.makedirs(output_dir, exist_ok=True)
        
        animated_videos = []
        failed_animations = []
        total_videos = len(audio_files)
        
        for i, audio_path in enumerate(audio_files):
            base_name = os.path.splitext(os.path.basename(audio_path))[0]
            output_path = os.path.join(output_dir, f"{base_name}_animated.mp4")
            
            # Extract slide number from filename
            slide_num = base_name.replace('slide_', '').lstrip('0') or '1'
            
            # Update progress and status
            progress = 50 + (i * 25) // total_videos  # 50-75% range
            current_step = f"Animating face for slide {slide_num}/{total_videos}"
            print(f"\n🎭 PROCESSING SLIDE {slide_num}/{total_videos} - Animating face...")
            self.logger.info(f"Animating face for slide {slide_num}/{total_videos} (Progress: {progress}%)")
            
            # Update global status for frontend
            self.update_status(progress, current_step)
            
            if (os.path.exists(output_path) and 
                self.config.config['processing']['skip_existing'] and
                os.path.getsize(output_path) > 0):
                self.logger.info(f"SKIP: {base_name}, video already exists")
                animated_videos.append(output_path)
                continue
            
            # Try to animate face with retry logic
            max_retries = 2
            success = False
            
            for attempt in range(max_retries):
                current_step = f"Animating face for slide {slide_num}/{total_videos} (attempt {attempt + 1}/{max_retries})"
                self.logger.info(f"Animating face for {base_name} (attempt {attempt + 1}/{max_retries})")
                self.update_status(progress, current_step)
                
                if self.animate_face(audio_path, face_path, output_path):
                    animated_videos.append(output_path)
                    success = True
                    print(f"✅ FACE ANIMATED FOR SLIDE {slide_num}/{total_videos} - Moving to next slide")
                    self.logger.info(f"SUCCESS: Face animated for slide {slide_num}")
                    # Update status to show completion
                    current_step = f"Face animated for slide {slide_num}/{total_videos} - Moving to next slide"
                    self.update_status(progress, current_step)
                    break
                else:
                    self.logger.warning(f"Attempt {attempt + 1} failed for {base_name}")
                    if attempt < max_retries - 1:
                        time.sleep(5)  # Wait before retry
            
            if not success:
                self.logger.error(f"ERROR: Failed to animate face for {base_name} after {max_retries} attempts")
                failed_animations.append(base_name)
            else:
                # Small delay to ensure status update is sent
                time.sleep(0.5)
        
        if failed_animations:
            self.logger.warning(f"WARNING: Failed to animate faces for: {failed_animations}")
        
        self.logger.info(f"SUCCESS: Generated {len(animated_videos)} animated videos out of {len(audio_files)} audio files")
        return animated_videos

class VideoComposer:
    """Video composer using ffmpeg"""
    
    def __init__(self, config: EasyConfig):
        self.config = config
        self.logger = logging.getLogger(__name__)
        self.system_instance = None  # Will be set by the main system
    
    def set_system_instance(self, system_instance):
        """Set reference to main system for status updates"""
        self.system_instance = system_instance
    
    def update_status(self, progress: int, current_step: str):
        """Update the global processing status"""
        if self.system_instance:
            self.system_instance.processing_status['progress'] = progress
            self.system_instance.processing_status['current_step'] = current_step
            # Show status update prominently in terminal
            print(f"\n🔄 STATUS UPDATE: {progress}% - {current_step}")
            self.logger.info(f"Status updated: {progress}% - {current_step}")
    
    def compose_final_video_with_slides(self) -> bool:
        """Compose final video with slide backgrounds and animated avatars"""
        self.logger.info("Creating final video with slide backgrounds...")
        
        # Check if we have the required files
        video_dir = "video"
        slide_dir = "slide_images"
        output_dir = "output"
        
        # Ensure directories exist
        os.makedirs(video_dir, exist_ok=True)
        os.makedirs(slide_dir, exist_ok=True)
        os.makedirs(output_dir, exist_ok=True)
        
        # Find animated videos
        animated_videos = []
        if os.path.exists(video_dir):
            for file in os.listdir(video_dir):
                if file.endswith("_animated.mp4"):
                    file_path = os.path.join(video_dir, file)
                    if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
                        animated_videos.append(file_path)
        
        # Find slide images
        slide_images = []
        if os.path.exists(slide_dir):
            for file in os.listdir(slide_dir):
                if file.endswith(".png"):
                    file_path = os.path.join(slide_dir, file)
                    if os.path.exists(file_path) and os.path.getsize(file_path) > 0:
                        slide_images.append(file_path)
        
        self.logger.info(f"Found {len(animated_videos)} animated videos")
        self.logger.info(f"Found {len(slide_images)} slide images")
        
        if not animated_videos:
            self.logger.error("ERROR: No animated videos found")
            return False
        
        if not slide_images:
            self.logger.error("ERROR: No slide images found")
            return False
        
        # Sort files to match them properly
        animated_videos.sort()
        slide_images.sort()
        
        # Verify we have matching pairs
        if len(animated_videos) != len(slide_images):
            self.logger.warning(f"WARNING: Mismatch - {len(animated_videos)} videos vs {len(slide_images)} slides")
            # Use the minimum count
            min_count = min(len(animated_videos), len(slide_images))
            animated_videos = animated_videos[:min_count]
            slide_images = slide_images[:min_count]
        
        try:
            # Create temporary directory
            temp_dir = "temp"
            os.makedirs(temp_dir, exist_ok=True)
            
            # Process each slide
            processed_videos = []
            total_slides = len(animated_videos)
            
            for i, (video_path, slide_path) in enumerate(zip(animated_videos, slide_images)):
                slide_num = i + 1
                progress = 75 + (i * 20) // total_slides  # 75-95% range
                current_step = f"Composing video for slide {slide_num}/{total_slides}"
                self.logger.info(f"Processing slide {slide_num}/{total_slides} (Progress: {progress}%)")
                
                # Update global status for frontend
                self.update_status(progress, current_step)
                
                if not os.path.exists(video_path) or not os.path.exists(slide_path):
                    self.logger.warning(f"WARNING: Missing files for slide {slide_num}, skipping")
                    continue
                
                # Create a composite video with slide background and avatar
                output_path = os.path.join(temp_dir, f"slide_{i+1:03d}_composite.mp4")
                
                # Use ffmpeg to overlay the animated avatar on the slide image
                # Position avatar in bottom right corner with proper size for lip-sync visibility
                cmd = [
                    "ffmpeg", "-y",
                    "-loop", "1", "-i", slide_path,  # Slide image as background
                    "-i", video_path,  # Animated avatar video
                    "-filter_complex", 
                    f"[0:v]scale=1920:1080[bg];[1:v]scale=320:320[avatar];[bg][avatar]overlay=1600:760[out]",  # Smaller avatar, bottom right corner
                    "-map", "[out]",
                    "-map", "1:a",  # Use audio from avatar video
                    "-c:v", "libx264",
                    "-c:a", "aac",
                    "-shortest",
                    "-t", "30",  # Limit to 30 seconds max per slide
                    output_path
                ]
                
                self.logger.info(f"  Creating composite video for slide {i+1}...")
                self.logger.info(f"  Command: {' '.join(cmd)}")
                
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
                
                if result.returncode == 0 and os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                    processed_videos.append(output_path)
                    self.logger.info(f"  SUCCESS: Created composite video for slide {i+1}")
                else:
                    self.logger.error(f"  ERROR: Failed to create composite video for slide {i+1}")
                    self.logger.error(f"  Return code: {result.returncode}")
                    self.logger.error(f"  Error: {result.stderr}")
                    self.logger.error(f"  Output: {result.stdout}")
            
            if not processed_videos:
                self.logger.error("ERROR: No composite videos created")
                return False
            
            # Concatenate all composite videos
            self.logger.info("Concatenating all slides...")
            self.update_status(95, "Concatenating final video...")
            
            # Create file list for ffmpeg
            file_list_path = os.path.join(temp_dir, "video_list.txt")
            with open(file_list_path, 'w') as f:
                for video_path in processed_videos:
                    f.write(f"file '{os.path.abspath(video_path)}'\n")
            
            # Concatenate videos
            final_output = os.path.join(output_dir, "final_presentation_with_slides.mp4")
            concat_cmd = [
                "ffmpeg", "-y",
                "-f", "concat",
                "-safe", "0",
                "-i", file_list_path,
                "-c", "copy",
                final_output
            ]
            
            self.logger.info("Concatenating videos...")
            result = subprocess.run(concat_cmd, capture_output=True, text=True)
            
            if result.returncode == 0:
                self.logger.info(f"SUCCESS: Final video created: {final_output}")
                return True
            else:
                self.logger.error(f"ERROR: Failed to concatenate videos: {result.stderr}")
                return False
        
        except Exception as e:
            self.logger.error(f"ERROR: Video composition failed: {e}")
            return False

class EasySpeakerAvatarSystem:
    """Easy unified system class"""
    
    def __init__(self):
        self.config = EasyConfig()
        self.logger = logging.getLogger(__name__)
        
        # Initialize processors
        self.ppt_processor = PowerPointProcessor(self.config)
        self.slide_renderer = SlideRenderer()
        self.tts_processor = TTSProcessor(self.config)
        self.wav2lip_processor = Wav2LipProcessor(self.config)
        self.video_composer = VideoComposer(self.config)
        
        # Set system instance references for status updates
        self.tts_processor.set_system_instance(self)
        self.wav2lip_processor.set_system_instance(self)
        self.video_composer.set_system_instance(self)
        
        # Create directories
        self._create_directories()
        
        # Processing status
        self.processing_status = {
            'status': 'idle',
            'progress': 0,
            'current_step': '',
            'error': None
        }
    
    def _create_directories(self):
        """Create all necessary directories"""
        for dir_name in self.config.config['directories'].values():
            os.makedirs(dir_name, exist_ok=True)
        self.logger.info("SUCCESS: All directories created")
    
    def set_gemini_api_key(self, api_key: str):
        """Set Gemini API key for enhanced explanations"""
        self.config.config['gemini']['api_key'] = api_key
        self.config.config['gemini']['enabled'] = bool(api_key.strip())
        self.logger.info(f"Gemini API {'enabled' if self.config.config['gemini']['enabled'] else 'disabled'}")
    
    def set_tts_voice(self, voice_name: str):
        """Set TTS voice for audio generation"""
        self.config.config['tts']['voice'] = voice_name
        self.tts_processor.tts_manager.set_voice(voice_name)
        self.logger.info(f"TTS voice set to: {voice_name}")
    
    def process_presentation_async(self, pptx_path: str, face_path: str):
        """Process presentation asynchronously"""
        def process():
            try:
                self.processing_status['status'] = 'processing'
                self.processing_status['progress'] = 0
                self.processing_status['current_step'] = 'Starting processing...'
                self.processing_status['error'] = None
                
                # Step 1: Extract text from file
                self.processing_status['progress'] = 5
                self.processing_status['current_step'] = 'Extracting text from slides...'
                print(f"\n📋 STEP 1: Extracting text from slides...")
                self.logger.info("Step 1: Extracting text from slides...")
                slides_data = self.ppt_processor.extract_slides_from_file(pptx_path)
                if not slides_data:
                    raise Exception("Failed to extract slides")
                print(f"✅ STEP 1 COMPLETED: Extracted {len(slides_data)} slides")
                self.logger.info(f"Step 1 completed: Extracted {len(slides_data)} slides")
                
                # Step 2: Render slide images
                self.processing_status['progress'] = 15
                self.processing_status['current_step'] = 'Rendering slide images...'
                print(f"\n🖼️ STEP 2: Rendering slide images...")
                self.logger.info("Step 2: Rendering slide images...")
                slide_images = self.slide_renderer.render_presentation(
                    pptx_path, 
                    self.config.config['directories']['slide_images']
                )
                if not slide_images:
                    raise Exception("Failed to render slide images")
                print(f"✅ STEP 2 COMPLETED: Rendered {len(slide_images)} slide images")
                self.logger.info(f"Step 2 completed: Rendered {len(slide_images)} slide images")
                
                # Step 3: Generate audio
                self.processing_status['progress'] = 25
                self.processing_status['current_step'] = 'Generating audio...'
                print(f"\n🎵 STEP 3: Generating audio for all slides...")
                self.logger.info("Step 3: Generating audio...")
                audio_files = self.tts_processor.generate_audio_batch(
                    slides_data, 
                    self.config.config['directories']['audio']
                )
                if not audio_files:
                    raise Exception("Failed to generate audio")
                print(f"✅ STEP 3 COMPLETED: Generated {len(audio_files)} audio files")
                self.logger.info(f"Step 3 completed: Generated {len(audio_files)} audio files")
                
                # Step 4: Animate faces
                self.processing_status['progress'] = 50
                self.processing_status['current_step'] = 'Animating faces...'
                print(f"\n🎭 STEP 4: Animating faces for all slides...")
                self.logger.info("Step 4: Animating faces...")
                animated_videos = self.wav2lip_processor.animate_faces_batch(
                    audio_files,
                    face_path,
                    self.config.config['directories']['video']
                )
                if not animated_videos:
                    raise Exception("Failed to animate faces")
                print(f"✅ STEP 4 COMPLETED: Generated {len(animated_videos)} animated videos")
                self.logger.info(f"Step 4 completed: Generated {len(animated_videos)} animated videos")
                
                # Step 5: Compose final video
                self.processing_status['progress'] = 75
                self.processing_status['current_step'] = 'Composing final video...'
                print(f"\n🎬 STEP 5: Composing final video...")
                self.logger.info("Step 5: Composing final video...")
                success = self.video_composer.compose_final_video_with_slides()
                if not success:
                    raise Exception("Failed to compose final video")
                print(f"✅ STEP 5 COMPLETED: Final video composed successfully")
                self.logger.info("Step 5 completed: Final video composed successfully")
                
                # Complete
                self.processing_status['progress'] = 100
                self.processing_status['status'] = 'completed'
                self.processing_status['current_step'] = 'Processing complete!'
                self.processing_status['error'] = None
                
                self.logger.info("SUCCESS: Presentation processing completed!")
                
            except Exception as e:
                self.processing_status['status'] = 'error'
                self.processing_status['error'] = str(e)
                self.processing_status['current_step'] = f'Error: {str(e)}'
                self.logger.error(f"ERROR: Processing failed: {e}")
        
        # Start processing in background thread
        thread = threading.Thread(target=process)
        thread.daemon = True
        thread.start()
    
    def get_processing_status(self):
        """Get current processing status"""
        return self.processing_status
    
    def get_slides_data(self, file_path: str) -> List[Dict[str, str]]:
        """Get slides data for frontend display"""
        return self.ppt_processor.extract_slides_from_file(file_path)

# Global system instance
system = EasySpeakerAvatarSystem()

# Flask routes
@app.route('/')
def index():
    """Serve the main frontend page"""
    return send_from_directory('frontend', 'index.html')

@app.route('/<path:filename>')
def serve_static(filename):
    """Serve static files"""
    return send_from_directory('frontend', filename)

@app.route('/api/process', methods=['POST'])
def process_presentation():
    """Process presentation and return slides data"""
    try:
        # Check if files are uploaded
        if 'pptx_file' not in request.files:
            return jsonify({'error': 'No presentation file uploaded'}), 400
        
        if 'face_file' not in request.files:
            return jsonify({'error': 'No face image uploaded'}), 400
        
        presentation_file = request.files['pptx_file']
        face_file = request.files['face_file']
        
        if presentation_file.filename == '' or face_file.filename == '':
            return jsonify({'error': 'No files selected'}), 400
        
        # Get file extension and save with appropriate name
        file_ext = os.path.splitext(presentation_file.filename)[1].lower()
        if file_ext not in ['.pptx', '.pdf']:
            return jsonify({'error': 'Unsupported file format. Please upload .pptx or .pdf files.'}), 400
        
        # Save uploaded files
        presentation_path = os.path.join('uploads', f'presentation{file_ext}')
        face_path = os.path.join('uploads', 'face.png')
        
        os.makedirs('uploads', exist_ok=True)
        presentation_file.save(presentation_path)
        face_file.save(face_path)
        
        # Gemini API key is now built-in and enabled by default
        
        # Get selected voice if provided
        selected_voice = request.form.get('selected_voice', '')
        if selected_voice:
            system.set_tts_voice(selected_voice)
        
        # Extract slides data for immediate response
        slides_data = system.get_slides_data(presentation_path)
        
        # Start background processing
        system.process_presentation_async(presentation_path, face_path)
        
        # Return slides data for frontend display
        return jsonify({
            'success': True,
            'slides': [
                {
                    'slide_number': slide['slide_number'],
                    'title': slide['title'],
                    'content': slide['content'],
                    'speaker_notes': slide['speaker_notes'],
                    'narration': slide['narration_text']
                }
                for slide in slides_data
            ]
        })
        
    except Exception as e:
        logger.error(f"Error processing presentation: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/status')
def get_status():
    """Get processing status"""
    return jsonify(system.get_processing_status())

@app.route('/api/voices')
def get_voices():
    """Get available TTS voices (filtered to show only the best ones)"""
    try:
        voices = []
        if system.tts_processor.tts_manager.is_available():
            import pyttsx3
            engine = pyttsx3.init()
            system_voices = engine.getProperty('voices')
            
            # Filter voices to show only the best ones
            filtered_voices = []
            seen_names = set()
            
            for voice in system_voices:
                voice_name = voice.name
                # Skip duplicates and low-quality voices
                if (voice_name not in seen_names and 
                    not any(skip in voice_name.lower() for skip in ['mobile', 'telephone', 'whisper', 'robot'])):
                    
                    # Prioritize Microsoft voices and common languages
                    priority = 0
                    if 'microsoft' in voice_name.lower():
                        priority = 3
                    elif any(lang in voice_name.lower() for lang in ['english', 'en-']):
                        priority = 2
                    elif any(lang in voice_name.lower() for lang in ['us', 'uk', 'australia']):
                        priority = 1
                    
                    filtered_voices.append({
                        'name': voice_name,
                        'id': voice.id,
                        'language': getattr(voice, 'languages', ['en'])[0] if hasattr(voice, 'languages') else 'en',
                        'priority': priority
                    })
                    seen_names.add(voice_name)
            
            # Sort by priority (highest first) and then by name
            filtered_voices.sort(key=lambda x: (-x['priority'], x['name']))
            
            # Limit to top 20 voices
            voices = filtered_voices[:20]
            
        return jsonify({'voices': voices})
    except Exception as e:
        logger.error(f"Error getting voices: {e}")
        return jsonify({'voices': []})

@app.route('/api/preview-voice', methods=['POST'])
def preview_voice():
    """Generate a voice preview and return download URL"""
    try:
        data = request.get_json()
        voice_name = data.get('voice_name', '')
        
        if not voice_name:
            return jsonify({'error': 'No voice name provided'}), 400
        
        logger.info(f"Generating voice preview for: {voice_name}")
        
        # Create preview directory
        preview_dir = "previews"
        os.makedirs(preview_dir, exist_ok=True)
        
        # Create a safe filename
        safe_voice_name = "".join(c for c in voice_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
        safe_voice_name = safe_voice_name.replace(' ', '_')
        preview_filename = f"preview_{safe_voice_name}.wav"
        preview_path = os.path.join(preview_dir, preview_filename)
        
        # Create a new TTS manager instance for preview to avoid conflicts
        try:
            import pyttsx3
            preview_engine = pyttsx3.init()
            preview_engine.setProperty('rate', 150)
            preview_engine.setProperty('volume', 0.9)
            
            # Set the voice
            voices = preview_engine.getProperty('voices')
            voice_found = False
            for voice in voices:
                if voice.name == voice_name:
                    preview_engine.setProperty('voice', voice.id)
                    voice_found = True
                    break
            
            if not voice_found:
                logger.warning(f"Voice {voice_name} not found, using default")
            
            # Generate preview audio
            preview_text = "This is a voice preview."
            preview_engine.save_to_file(preview_text, preview_path)
            preview_engine.runAndWait()
            
            # Wait for file to be written
            time.sleep(1)
            
            if os.path.exists(preview_path) and os.path.getsize(preview_path) > 0:
                logger.info(f"SUCCESS: Generated voice preview for {voice_name}")
                return jsonify({
                    'success': True,
                    'download_url': f'/api/download-preview/{preview_filename}',
                    'voice_name': voice_name
                })
            else:
                logger.error(f"Failed to generate preview audio file")
                return jsonify({'error': 'Failed to generate preview audio'}), 500
                
        except Exception as e:
            logger.error(f"Error in voice preview generation: {e}")
            return jsonify({'error': f'Voice preview generation failed: {str(e)}'}), 500
            
    except Exception as e:
        logger.error(f"Error generating voice preview: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/download-preview/<filename>')
def download_preview(filename):
    """Download a voice preview file"""
    try:
        preview_dir = "previews"
        file_path = os.path.join(preview_dir, filename)
        
        if os.path.exists(file_path):
            return send_from_directory(preview_dir, filename, as_attachment=True)
        else:
            return jsonify({'error': 'Preview file not found'}), 404
            
    except Exception as e:
        logger.error(f"Error downloading preview: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/download')
def download_video():
    """Download the final video"""
    video_path = os.path.join('output', 'final_presentation_with_slides.mp4')
    if os.path.exists(video_path):
        return send_from_directory('output', 'final_presentation_with_slides.mp4', as_attachment=True)
    else:
        return jsonify({'error': 'Video not ready yet'}), 404

@app.route('/api/cleanup', methods=['POST'])
def cleanup_files():
    """Clean up generated files"""
    try:
        # Clean up directories
        directories_to_clean = ['audio', 'video', 'output', 'slide_images', 'temp', 'uploads', 'previews']
        
        for directory in directories_to_clean:
            if os.path.exists(directory):
                for file in os.listdir(directory):
                    file_path = os.path.join(directory, file)
                    if os.path.isfile(file_path):
                        os.remove(file_path)
        
        # Reset processing status
        system.processing_status = {
            'status': 'idle',
            'progress': 0,
            'current_step': '',
            'error': None
        }
        
        # Reset TTS engine
        system.tts_processor.tts_manager.reset_engine()
        
        return jsonify({'success': True, 'message': 'Files cleaned up and system reset successfully'})
        
    except Exception as e:
        logger.error(f"Error cleaning up files: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/reset', methods=['POST'])
def reset_system():
    """Reset the system to clear stuck states"""
    try:
        # Reset processing status
        system.processing_status = {
            'status': 'idle',
            'progress': 0,
            'current_step': '',
            'error': None
        }
        
        # Reset TTS engine
        system.tts_processor.tts_manager.reset_engine()
        
        logger.info("System reset successfully")
        return jsonify({'success': True, 'message': 'System reset successfully'})
        
    except Exception as e:
        logger.error(f"Error resetting system: {e}")
        return jsonify({'error': str(e)}), 500

def open_browser():
    """Open the web browser automatically"""
    import webbrowser
    import time
    time.sleep(2)  # Wait for server to start
    webbrowser.open('http://localhost:5000')

def main():
    """Main function to run the unified system"""
    print("🎬 Easy AI Speaker Avatar System")
    print("=" * 50)
    print("🚀 Starting unified system...")
    print("🌐 Frontend will be available at: http://localhost:5000")
    print("📡 API endpoints:")
    print("  POST /api/process - Process presentation")
    print("  GET  /api/status - Get processing status")
    print("  GET  /api/download - Download final video")
    print("  POST /api/cleanup - Clean up files")
    print("=" * 50)
    
    print("\n🎉 System ready! Opening web interface...")
    print("📝 Upload your PowerPoint file and face image to get started!")
    
    # Open browser in a separate thread
    browser_thread = threading.Thread(target=open_browser)
    browser_thread.daemon = True
    browser_thread.start()
    
    # Start Flask server
    print("🌐 Server starting...")
    print("🤖 Gemini AI enabled for enhanced explanations!")
    app.run(host='0.0.0.0', port=5000, debug=False, threaded=True, use_reloader=False)

if __name__ == "__main__":
    main()
