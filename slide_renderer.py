# Slide Renderer - Enhanced PowerPoint slide image extraction
# This module provides better slide rendering capabilities

import os
import sys
from pathlib import Path
from typing import List, Tuple
import logging

from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw, ImageFont
import io

class SlideRenderer:
    """Enhanced slide renderer for PowerPoint presentations"""
    
    def __init__(self, width: int = 1920, height: int = 1080):
        self.width = width
        self.height = height
        self.logger = logging.getLogger(__name__)
        
        # Try to load a system font
        try:
            self.title_font = ImageFont.truetype("arial.ttf", 48)
            self.content_font = ImageFont.truetype("arial.ttf", 32)
            self.notes_font = ImageFont.truetype("arial.ttf", 24)
        except:
            # Fallback to default font
            self.title_font = ImageFont.load_default()
            self.content_font = ImageFont.load_default()
            self.notes_font = ImageFont.load_default()
    
    def render_slide(self, slide, slide_number: int) -> Image.Image:
        """Render a single slide as an image"""
        # Create a new image with white background
        img = Image.new('RGB', (self.width, self.height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Define layout areas
        margin = 50
        title_area = (margin, margin, self.width - margin, 150)
        content_area = (margin, 180, self.width - margin, self.height - 200)
        notes_area = (margin, self.height - 150, self.width - margin, self.height - margin)
        
        # Draw slide number
        draw.text((self.width - 100, 20), f"Slide {slide_number}", 
                 fill='gray', font=self.content_font)
        
        # Extract and render title
        if slide.shapes.title and slide.shapes.title.text.strip():
            title_text = slide.shapes.title.text.strip()
            self._draw_wrapped_text(draw, title_text, title_area, self.title_font, 'black')
        
        # Extract and render content
        content_text = self._extract_slide_content(slide)
        if content_text:
            self._draw_wrapped_text(draw, content_text, content_area, self.content_font, 'black')
        
        # Extract and render speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                # Draw a separator line
                draw.line([margin, self.height - 160, self.width - margin, self.height - 160], 
                         fill='lightgray', width=2)
                self._draw_wrapped_text(draw, f"Notes: {notes_text}", notes_area, 
                                      self.notes_font, 'darkgray')
        
        return img
    
    def _extract_slide_content(self, slide) -> str:
        """Extract all text content from slide shapes"""
        content_parts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Skip the title shape
                if shape != slide.shapes.title:
                    content_parts.append(shape.text.strip())
        
        return '\n'.join(content_parts)
    
    def _draw_wrapped_text(self, draw, text: str, area: Tuple, font, color: str):
        """Draw text with word wrapping within the specified area"""
        x, y, max_x, max_y = area
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            test_line = ' '.join(current_line + [word])
            bbox = draw.textbbox((0, 0), test_line, font=font)
            text_width = bbox[2] - bbox[0]
            
            if text_width <= max_x - x:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                    current_line = [word]
                else:
                    # Single word is too long, add it anyway
                    lines.append(word)
        
        if current_line:
            lines.append(' '.join(current_line))
        
        # Draw the lines
        line_height = 40  # Approximate line height
        for i, line in enumerate(lines):
            line_y = y + (i * line_height)
            if line_y + line_height <= max_y:
                draw.text((x, line_y), line, fill=color, font=font)
    
    def render_presentation(self, pptx_path: str, output_dir: str) -> List[str]:
        """Render all slides from a presentation"""
        self.logger.info(f"🖼️ Rendering slides from {pptx_path}")
        
        try:
            prs = Presentation(pptx_path)
            slide_images = []
            
            for i, slide in enumerate(prs.slides):
                slide_image = self.render_slide(slide, i + 1)
                output_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
                slide_image.save(output_path)
                slide_images.append(output_path)
                self.logger.info(f"   -> Rendered slide {i+1}")
            
            self.logger.info(f"✅ Rendered {len(slide_images)} slides")
            return slide_images
            
        except Exception as e:
            self.logger.error(f"❌ Error rendering slides: {e}")
            return []

def main():
    """Test the slide renderer"""
    import argparse
    
    parser = argparse.ArgumentParser(description='Render PowerPoint slides as images')
    parser.add_argument('pptx_path', help='Path to PowerPoint file')
    parser.add_argument('output_dir', help='Output directory for slide images')
    parser.add_argument('--width', type=int, default=1920, help='Image width')
    parser.add_argument('--height', type=int, default=1080, help='Image height')
    
    args = parser.parse_args()
    
    # Setup logging
    logging.basicConfig(level=logging.INFO)
    
    # Create renderer and process
    renderer = SlideRenderer(args.width, args.height)
    os.makedirs(args.output_dir, exist_ok=True)
    
    slide_images = renderer.render_presentation(args.pptx_path, args.output_dir)
    print(f"Rendered {len(slide_images)} slides to {args.output_dir}")

if __name__ == "__main__":
    main()
